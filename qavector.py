import streamlit as st
from openai import OpenAI
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
import mammoth
import requests
import openpyxl
from qdrant_client import QdrantClient

# Função para ler arquivos PDF
def read_pdf(file):
    try:
        document = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in document:
            text += page.get_text()
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PDF: {e}")
        return ""

# Função para ler arquivos DOCX
def read_docx(file):
    try:
        document = Document(file)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOCX: {e}")
        return ""

# Função para ler arquivos DOC usando mammoth
def read_doc(file):
    try:
        result = mammoth.extract_raw_text(file)
        return result.value
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOC: {e}")
        return ""

# Função para ler arquivos PPT e PPTX
def read_ppt_pptx(file):
    try:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PPT/PPTX: {e}")
        return ""

# Função para ler arquivos TXT e MD com detecção automática de codificação
def read_txt_md(file, encoding="utf-8"):
    try:
        # Detectar a codificação do arquivo
        detected_encoding = chardet.detect(file.read())["encoding"]

        # Reabrir o arquivo com a codificação detectada
        file.seek(0)
        text = file.read().decode(encoding or detected_encoding)

        return text

    except Exception as e:
        st.error(f"Erro ao ler arquivo TXT/MD: {e}")
        return ""

def read_xls(file_path):
    try:
        # Abrir o arquivo XLS
        workbook = openpyxl.load_workbook(file_path)

        # Acessar a primeira planilha do arquivo
        worksheet = workbook.active

        # Criar um dicionário para armazenar os dados
        data = {}

        # Iterar sobre as linhas da planilha
        for row in worksheet.iter_rows():
            # Iterar sobre as células da linha atual
            for cell in row:
                # Verificar se a célula contém dados
                if cell.value:
                    # Adicionar os dados ao dicionário
                    data[cell.coordinate] = cell.value

        # Fechar o arquivo
        workbook.close()

        return data

    except FileNotFoundError:
        #print(f"Arquivo '{file_path}' não encontrado.")
        st.write(f"Arquivo '{file_path}' não encontrado.")
        return None

    except Exception as e:
        #print(f"Erro ao ler o arquivo XLS: {e}")
        sr.write(f"Erro ao ler o arquivo XLS: {e}")
        return None

# Função para converter o arquido carregado para texto
def trata_arquivo (uploaded_file):
    #st.write("Tipo arquivo carregado",uploaded_file.type)
    # Process the uploaded file based on its type
    if uploaded_file.type == "application/pdf":
        document_text = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document_text = read_docx(uploaded_file)
    elif uploaded_file.type == "application/msword":
        document_text = read_doc(uploaded_file)
    elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        document_text = read_ppt_pptx(uploaded_file)
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/xls","application/xlsx", "application/xlsm","application/xltx","application/xltm"]: 
        document_text = read_xls(uploaded_file)
    else:
        document_text = read_txt_md(uploaded_file)
    return document_text

# Função para limpar o conteúdo da variável "question"
def clear_question():
    try:
        del question
    except NameError:
        pass

# Criar uma caixa de entrada de texto
def get_question():
    clear_question()
    return st.text_input(
        "Faça um questionamento",
        placeholder="Por exemplo: Pode fornecer um sumário?",
        disabled=not uploaded_file,
    )

# Streamlit UI
#st.title("📝 Carregue o Documento")
st.markdown("## 📝 Carregue o Documento")

# Create an OpenAI client.
#client = OpenAI(api_key=retrieved_secret.value) versão AZURE
cliente = OpenAI(api_key=st.secrets.OPENAI_KEY)

# Criar um objeto Qdrant para persistir os vetores
#qdrant_client = Qdrant("http://localhost:6333")
qdrant_client = QdrantClient(":memory:")

# Let the user upload a file via `st.file_uploader`.
uploaded_files = st.file_uploader("Selecione os Documentos a serem analisados!", 
                                  type=("pdf", "docx", "doc", "ppt", "pptx", "txt", "md","xls","xlsx","xlsm","xltx","xltm"),
                                  accept_multiple_files=True)
for uploaded_file in uploaded_files:
     # Imprimir informações sobre o arquivo
    st.write("Nome do arquivo:", uploaded_file.name)
    st.write("Tipo de conteúdo:", uploaded_file.type)
    st.write("Tamanho do arquivo:", uploaded_file.size, "bytes")
    document = trata_arquivo(uploaded_file)
    # Criar embeddings para o arquivo
    chunk_size = 1024
    embeddings = []
    with document.open() as file:
        while True:
            chunk = file.read(chunk_size)
            if not chunk:
                break
    response = openai.Embedding.create(input=chunk)
    embedding = response["data"]["embedding"]
    qdrant_client.add(embedding)
    embeddings.append(embedding)
    # Adicionar os embeddings do arquivo à lista de todos os embeddings
    all_embeddings.extend(embeddings)

# Obter a pergunta do usuário após o processamento de todos os arquivos
question = st.text_input(
    "Faça um questionamento após a análise de todos os documentos",
    placeholder="Por exemplo: Pode fornecer um sumário?",
)

# Obter a pergunta do usuário após o processamento de todos os arquivos
question = st.text_input(
    "Faça um questionamento após a análise de todos os documentos",
    placeholder="Por exemplo: Pode fornecer um sumário?",
)

if question:
    # Encontre o embedding mais próximo à pergunta no Qdrant
    closest_embedding = qdrant_client.search(question, in_memory=True)[0]

    # Recupere o índice do embedding mais próximo
    closest_index = closest_embedding["id"]

    # Recupere o texto original do índice mais próximo
    original_text = all_embeddings[closest_index]

    # Gere uma resposta à pergunta usando a API da OpenAI
    completion = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Pergunta: {question}\nContexto: {original_text}\nResposta:",
        max_tokens= 50,
        n= 1,
        stop= ["\n"],
    )

    # Exibir a resposta na interface Streamlit
    st.write(f"Pergunta: {question}\nResposta: {completion.choices[0]['text']}")
