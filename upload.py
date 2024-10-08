import streamlit as st
from openai import OpenAI
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_qdrant import Qdrant
from langchain_text_splitters import TokenTextSplitter
import mammoth
import requests
import openpyxl
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# Função para ler arquivos PDF
def read_pdf(file):
    try:
        document = fitz.open(stream=file.read(), filetype="pdf")
        text = ""
        for page in document:
            text += page.get_text()
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PDF: {e}")
        return ""

# Função para ler arquivos DOCX
def read_docx(file):
    try:
        document = Document(file)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOCX: {e}")
        return ""

# Função para ler arquivos DOC usando mammoth
def read_doc(file):
    try:
        result = mammoth.extract_raw_text(file)
        return result.value
    except Exception as e:
        st.error(f"Erro ao ler arquivo DOC: {e}")
        return ""

# Função para ler arquivos PPT e PPTX
def read_ppt_pptx(file):
    try:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Erro ao ler arquivo PPT/PPTX: {e}")
        return ""

# Função para ler arquivos TXT e MD com detecção automática de codificação
def read_txt_md(file, encoding="utf-8"):
    try:
        # Detectar a codificação do arquivo
        detected_encoding = chardet.detect(file.read())["encoding"]

        # Reabrir o arquivo com a codificação detectada
        file.seek(0)
        text = file.read().decode(encoding or detected_encoding)

        return text

    except Exception as e:
        st.error(f"Erro ao ler arquivo TXT/MD: {e}")
        return ""

def read_xls(file_path):
    try:
        # Abrir o arquivo XLS
        workbook = openpyxl.load_workbook(file_path)

        # Acessar a primeira planilha do arquivo
        worksheet = workbook.active

        # Criar um dicionário para armazenar os dados
        data = {}

        # Iterar sobre as linhas da planilha
        for row in worksheet.iter_rows():
            # Iterar sobre as células da linha atual
            for cell in row:
                # Verificar se a célula contém dados
                if cell.value:
                    # Adicionar os dados ao dicionário
                    data[cell.coordinate] = cell.value

        # Fechar o arquivo
        workbook.close()

        return data

    except FileNotFoundError:
        #print(f"Arquivo '{file_path}' não encontrado.")
        st.write(f"Arquivo '{file_path}' não encontrado.")
        return None

    except Exception as e:
        #print(f"Erro ao ler o arquivo XLS: {e}")
        sr.write(f"Erro ao ler o arquivo XLS: {e}")
        return None

# Função para converter o arquido carregado para texto
def trata_arquivo (uploaded_file):
    #st.write("Tipo arquivo carregado",uploaded_file.type)
    # Process the uploaded file based on its type
    if uploaded_file.type == "application/pdf":
        #st.write("Lendo arquivo PDF",uploaded_file.name)
        document_text = read_pdf(uploaded_file)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        document_text = read_docx(uploaded_file)
    elif uploaded_file.type == "application/msword":
        document_text = read_doc(uploaded_file)
    elif uploaded_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
        document_text = read_ppt_pptx(uploaded_file)
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet","application/xls","application/xlsx", "application/xlsm","application/xltx","application/xltm"]: 
        document_text = read_xls(uploaded_file)
    else:
        document_text = read_txt_md(uploaded_file)
    return document_text

# Função para limpar o conteúdo da variável "question"
def clear_question():
    try:
        del question
    except NameError:
        pass

# Criar uma caixa de entrada de texto
def get_question():
    clear_question()
    return st.text_input(
        "Faça um questionamento",
        placeholder="Por exemplo: Pode fornecer um sumário?",
        disabled=not uploaded_file,
    )

# Streamlit UI
#st.title("📝 Selecione os documentos")
st.markdown("## 📝 Selecione os Documentos")
model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
model_kwargs = {'device': 'cpu'}
encode_kwargs = {'normalize_embeddings': True}
hf = HuggingFaceEmbeddings(
    model_name=model_name,
    model_kwargs=model_kwargs,
    encode_kwargs=encode_kwargs
)
# Let the user upload a file via `st.file_uploader`.
uploaded_files = st.file_uploader("Selecione os Documentos a serem analisados!", 
                                  type=("pdf", "docx", "doc", "ppt", "pptx", "txt", "md","xls","xlsx","xlsm","xltx","xltm"),
                                  accept_multiple_files=True)
for uploaded_file in uploaded_files:
     # Imprimir informações sobre o arquivo
    st.write("Nome do arquivo:", uploaded_file.name)
    st.write("Tipo de conteúdo:", uploaded_file.type)
    st.write("Tamanho do arquivo:", uploaded_file.size, "bytes")
    document = trata_arquivo(uploaded_file)
    #
    client = QdrantClient(path="qdrant/")
    collection_name = "MyCollection"
    if client.collection_exists(collection_name):
        client.delete_collection(collection_name)
    client.create_collection(collection_name,vectors_config=VectorParams(size=1024, distance=Distance.DOT))
    qdrant = Qdrant(client, collection_name, hf)
    st.writing("Indexing...")
    # Chunking
    text_splitter = TokenTextSplitter(chunk_size=1024, chunk_overlap=102)
    texts = text_splitter.split_text(document)
    metadata = []
    for i in range(0,len(texts)):
        metadata.append({"path":file})
        #qdrant.add_texts(texts,metadatas=metadata)
    st.write("numero de chunks: ",len(texts))
